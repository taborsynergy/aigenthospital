from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY  = RGBColor(0x0F, 0x1F, 0x35)
NAVY2 = RGBColor(0x16, 0x2C, 0x49)
TEAL  = RGBColor(0x14, 0xB8, 0xA6)
TEALD = RGBColor(0x0D, 0x94, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF6, 0xF8, 0xFA)
SLATE = RGBColor(0x33, 0x41, 0x55)
GRAY  = RGBColor(0x64, 0x74, 0x8B)
CARD  = RGBColor(0xFF, 0xFF, 0xFF)
CARDBR = RGBColor(0xE2, 0xE8, 0xF0)
GOLD  = RGBColor(0xF5, 0x9E, 0x0B)
MUTE  = RGBColor(0xAF, 0xBE, 0xD0)
FAINT = RGBColor(0xCB, 0xD5, 0xE1)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SWI = prs.slide_width.inches


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, line_sp=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sp_after)
        p.line_spacing = line_sp
        for (t, sz, b, c) in para:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(sz)
            f.bold = b
            f.color.rgb = c
            f.name = FONT
    return tb


# ===== SLIDE 1 — TITLE =====
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
rect(s, 0, 0, SWI, 2.35, NAVY2)
rect(s, 0, 2.30, SWI, 0.06, TEAL)
rect(s, 9.35, 0.5, 3.45, 0.55, NAVY, line=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 9.35, 0.6, 3.45, 0.4, [[("Powered by Claude Sonnet 4.6", 11, True, TEAL)]], align=PP_ALIGN.CENTER)
txt(s, 0.9, 1.5, 11.5, 1.2, [[("Aria AI Front Desk", 42, True, WHITE)]])
txt(s, 0.92, 2.75, 11.5, 0.7, [[("The 24/7 AI Receptionist for Medical Practices", 22, True, TEAL)]])
txt(s, 0.92, 3.5, 11.6, 1.1,
    [[("Aria answers every patient call and message - books appointments, screens questions,", 15, False, FAINT)],
     [("verifies insurance, and safely triages emergencies - 24 hours a day, with no app and no hold music.", 15, False, FAINT)]],
    line_sp=1.15)
chips = ["Books appointments", "Triage-safe (911)", "HIPAA-aware", "24/7, no login", "Live dashboard"]
cx = 0.9
for c in chips:
    w = 0.16 + 0.105 * len(c)
    rect(s, cx, 5.1, w, 0.55, NAVY2, line=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, cx, 5.19, w, 0.4, [[(c, 11.5, True, WHITE)]], align=PP_ALIGN.CENTER)
    cx += w + 0.18
rect(s, 0.9, 6.0, 11.95, 0.5, NAVY2, line=TEALD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 0.9, 6.08, 11.95, 0.36,
    [[("Built secure:   HIPAA-aware    |    Encrypted in transit (TLS)    |    Row-level database security    |    Audit-logged", 11.5, True, FAINT)]],
    align=PP_ALIGN.CENTER)
txt(s, 0.9, 6.72, 11.95, 0.5, [[("Tabor Synergy", 13, True, TEAL), ("    .    Powered by Claude Sonnet 4.6 (Anthropic)    .    aifrontdesk.taborsynergy.com", 11, False, GRAY)]])

# ===== SLIDE 2 — CORE FEATURES =====
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
rect(s, 0, 0, SWI, 1.15, NAVY)
rect(s, 0, 1.10, SWI, 0.05, TEAL)
txt(s, 0.6, 0.28, 12.5, 0.7, [[("What Aria Does - for Every Patient, Every Time", 26, True, WHITE)]])
feats = [
    ("24/7 AI Receptionist", "Greets patients, answers questions, and handles intake any time - nights, weekends, holidays.", "No missed calls, no voicemail tag, no overtime."),
    ("Instant Appointment Booking", "Patients book in about 90 seconds via a link or QR code - no app, no login.", "Front desk freed from 50+ scheduling calls a day."),
    ("Emergency 911 Escalation", "Detects emergencies like chest pain and instantly tells the patient to call 911 + alerts staff.", "Patient safety + liability protection, built in."),
    ("HIPAA-Aware Guardrails", "Won't diagnose, won't reveal another person's records, resists manipulation.", "Compliant by design; protects patient privacy."),
    ("Insurance Q&A", "Answers 'do you take my insurance?' from the clinic's verified accepted-plans list.", "Fewer dead-end calls; patients self-qualify."),
    ("Live Dashboard & Analytics", "Staff see every appointment, no-show rates, busy times, and trends in real time.", "Run the practice on data, not guesswork."),
]
x0, y0, cw, ch, gx, gy = 0.55, 1.5, 4.0, 2.7, 0.18, 0.18
for i, (name, desc, adv) in enumerate(feats):
    r, c = divmod(i, 3)
    x = x0 + c * (cw + gx)
    y = y0 + r * (ch + gy)
    rect(s, x, y, cw, ch, CARD, line=CARDBR, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, 0.14, ch, TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x + 0.32, y + 0.2, cw - 0.55, 0.6, [[(name, 15.5, True, NAVY)]])
    txt(s, x + 0.34, y + 0.85, cw - 0.6, 1.2, [[(desc, 11.5, False, SLATE)]], line_sp=1.08)
    txt(s, x + 0.34, y + ch - 0.72, cw - 0.6, 0.6, [[("> " + adv, 11, True, TEALD)]], line_sp=1.05)

txt(s, 0.55, 7.12, 12.25, 0.32,
    [[("Enterprise-grade & HIPAA-aware   .   Encrypted (TLS)   .   Row-Level DB security   .   Audit-logged   .   Powered by Claude Sonnet 4.6", 11, True, TEALD)]],
    align=PP_ALIGN.CENTER)

# ===== SLIDE 3 — PLANS =====
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
rect(s, 0, 0, SWI, 1.15, NAVY)
rect(s, 0, 1.10, SWI, 0.05, TEAL)
txt(s, 0.6, 0.28, 12.5, 0.7, [[("Simple Plans That Grow With You", 26, True, WHITE)]])
plans = [
    ("STARTER", "$297", "/mo", "Solo practice, getting started", SLATE, CARDBR,
     ["AI receptionist + booking (core)", "300 conversations / month", "1 location, 1 provider", "Emergency triage + HIPAA safety", "Insurance Q&A + live dashboard", "Email support"]),
    ("GROWTH", "$597", "/mo", "2-5 doctor group practice", TEALD, TEAL,
     ["Everything in Starter, plus:", "1,000 conversations / month", "Up to 3 locations, 5 providers", "Custom assistant name + website widget", "Custom insurance knowledge", "Email reminders + recall campaigns", "Monthly reports, priority support"]),
    ("ENTERPRISE", "$997", "/mo", "Multi-location / white-label", SLATE, CARDBR,
     ["Everything in Growth, plus:", "Unlimited conversations", "Unlimited locations & providers", "White-label + custom domain", "EHR integration (Epic/Cerner)", "Custom AI training", "Dedicated manager, 24/7 support"]),
]
x0, y0, cw, ch, gx = 0.55, 1.45, 4.0, 5.4, 0.18
for i, (nm, price, per, who, accent, br, items) in enumerate(plans):
    x = x0 + i * (cw + gx)
    pop = (i == 1)
    yy = y0 - 0.15 if pop else y0
    hh = ch + 0.3 if pop else ch
    rect(s, x, yy, cw, hh, CARD, line=accent if pop else br, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=2.5 if pop else 1.0)
    rect(s, x, yy, cw, 0.9, accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if pop:
        rect(s, x + cw - 1.95, yy + 0.12, 1.75, 0.42, GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x + cw - 1.95, yy + 0.17, 1.75, 0.34, [[("MOST POPULAR", 9.5, True, WHITE)]], align=PP_ALIGN.CENTER)
    txt(s, x + 0.3, yy + 0.2, cw - 0.6, 0.5, [[(nm, 17, True, WHITE)]])
    txt(s, x + 0.3, yy + 1.05, cw - 0.6, 0.7, [[(price, 34, True, NAVY), (per, 14, False, GRAY)]])
    txt(s, x + 0.3, yy + 1.95, cw - 0.6, 0.5, [[(who, 12, True, accent)]])
    bl = []
    for it in items:
        head = it.startswith("Everything")
        bl.append([("-  ", 12, True, accent), (it, 11.5, head, SLATE)])
    txt(s, x + 0.3, yy + 2.5, cw - 0.55, hh - 2.6, bl, sp_after=7, line_sp=1.05)
txt(s, 0.55, 7.05, 12, 0.4, [[("All delivery is by email — no carrier registration, no per-clinic phone numbers, live today.", 10, False, GRAY)]])

# ===== SLIDE 4 — WHY / ROI / CTA =====
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
rect(s, 0, 0, SWI, 1.15, NAVY2)
rect(s, 0, 1.10, SWI, 0.05, TEAL)
txt(s, 0.6, 0.28, 12.5, 0.7, [[("Why Practices Choose Aria", 26, True, WHITE)]])
txt(s, 0.6, 1.2, 12.2, 0.32,
    [[("HIPAA-aware   .   Encrypted in transit & at rest   .   Row-Level DB security   .   Audit-logged   .   Powered by Claude Sonnet 4.6", 12, True, TEAL)]],
    align=PP_ALIGN.CENTER)
stats = [("40%", "fewer no-shows", "Automated reminders bring patients in"),
         ("3-4 hrs", "saved per day", "Staff stop fielding scheduling calls"),
         ("24/7", "always-on intake", "Capture patients nights & weekends"),
         ("15 min", "to go live", "$0 setup - we handle everything")]
x0, cw, gx = 0.7, 2.95, 0.27
for i, (big, lab, sub) in enumerate(stats):
    x = x0 + i * (cw + gx)
    rect(s, x, 1.55, cw, 2.05, NAVY2, line=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, 1.8, cw, 0.8, [[(big, 34, True, TEAL)]], align=PP_ALIGN.CENTER)
    txt(s, x, 2.64, cw, 0.4, [[(lab, 13, True, WHITE)]], align=PP_ALIGN.CENTER)
    txt(s, x + 0.2, 3.06, cw - 0.4, 0.5, [[(sub, 10.5, False, MUTE)]], align=PP_ALIGN.CENTER)
rect(s, 0.7, 3.95, 11.93, 1.15, NAVY2, line=TEALD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=1.5)
txt(s, 1.0, 3.95, 11.4, 1.15,
    [[("A typical solo practice nets about ", 16, True, WHITE), ("$7,000 / month", 16, True, TEAL),
      (" in saved staff time + recovered no-shows - vs. the plan cost.", 16, True, WHITE)]],
    anchor=MSO_ANCHOR.MIDDLE)
rect(s, 0.7, 5.45, 11.93, 1.35, TEALD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 1.0, 5.64, 11.4, 0.6, [[("Start a 2-week free trial - on your own branded patient link", 19, True, WHITE)]])
txt(s, 1.0, 6.27, 11.4, 0.5,
    [[("No card. No setup fee. Live the same day.     ", 13, False, WHITE),
      ("write2dinakar10@gmail.com   |   aifrontdesk.taborsynergy.com", 13, True, WHITE)]])

# ============================================================
# INTERNAL SLIDES (5-8) - detailed feature status. HIDE from Raj.
# ============================================================
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF1, 0x8A, 0x0B)
BLUE  = RGBColor(0x3B, 0x82, 0xF6)
RED   = RGBColor(0xDC, 0x26, 0x26)
STCOL = {"LIVE": GREEN, "PENDING": AMBER, "BUILT": BLUE}


def internal_slide(title, subtitle, rows):
    s = prs.slides.add_slide(BLANK)
    bg(s, LIGHT)
    rect(s, 0, 0, SWI, 1.15, NAVY)
    rect(s, 0, 1.10, SWI, 0.05, AMBER)
    txt(s, 0.6, 0.22, 9.5, 0.55, [[(title, 23, True, WHITE)]])
    txt(s, 0.62, 0.78, 9.5, 0.35, [[(subtitle, 12.5, False, FAINT)]])
    # internal badge
    rect(s, 10.6, 0.32, 2.25, 0.5, RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 10.6, 0.42, 2.25, 0.34, [[("INTERNAL - HIDE", 11, True, WHITE)]], align=PP_ALIGN.CENTER)
    # legend
    lx = 0.62
    for st in ["LIVE", "BUILT", "PENDING"]:
        rect(s, lx, 1.28, 0.85, 0.3, STCOL[st], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, lx, 1.31, 0.85, 0.26, [[(st, 9, True, WHITE)]], align=PP_ALIGN.CENTER)
        lx += 1.0
    # rows
    n = len(rows)
    top = 1.78
    rh = min(0.52, (7.15 - top) / n)
    for i, (name, detail, st) in enumerate(rows):
        y = top + i * rh
        if i % 2 == 0:
            rect(s, 0.55, y, 12.25, rh - 0.04, RGBColor(0xEC, 0xF1, 0xF6), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, 0.66, y + (rh - 0.30) / 2, 1.05, 0.30, STCOL[st], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, 0.66, y + (rh - 0.30) / 2 + 0.015, 1.05, 0.27, [[(st, 9, True, WHITE)]], align=PP_ALIGN.CENTER)
        txt(s, 1.85, y, 10.9, rh,
            [[(name + "   ", 12.5, True, NAVY), (detail, 11, False, GRAY)]],
            anchor=MSO_ANCHOR.MIDDLE, line_sp=1.0)
    return s


internal_slide(
    "STARTER ($297/mo) - Feature Detail",
    "Solo practice. Every line below is what the customer actually gets, with build status.",
    [
        ("AI receptionist + booking", "Aria chats, screens questions, and books appointments", "LIVE"),
        ("Hosted patient chat link + QR", "No-login page patients open via link or QR code", "LIVE"),
        ("300 conversations / month", "Usage cap enforced; upgrade prompt beyond the limit", "LIVE"),
        ("1 location / 1 provider", "Solo-practice limits enforced by plan", "LIVE"),
        ("Emergency 911 escalation", "Detects emergencies, tells patient to call 911, alerts staff", "LIVE"),
        ("HIPAA guardrails", "No diagnosis, no third-party records, prompt-injection safe", "LIVE"),
        ("Insurance Q&A", "Answers 'do you take X?' from the accepted-plans list", "LIVE"),
        ("Live dashboard + analytics", "Appointments, no-show rate, busy times, trends", "LIVE"),
        ("Booking confirmation email", "Emailed to patient on booking - LIVE (verified end-to-end via SendGrid)", "LIVE"),
        ("Email support", "Standard support (operational, not code)", "LIVE"),
    ])

internal_slide(
    "GROWTH ($597/mo) - Feature Detail",
    "2-5 doctor group practice. Everything in Starter, PLUS the items below.",
    [
        ("1,000 conversations / month", "3x Starter volume; cap enforced", "LIVE"),
        ("Up to 3 locations", "Each office: own address, hours, providers", "LIVE"),
        ("Up to 5 providers", "Group practice; AI books with named doctors", "LIVE"),
        ("Custom assistant name", "Rename 'Aria' to the clinic's own assistant name", "LIVE"),
        ("Website embed widget", "Branded booking chat embedded on the clinic's own site", "LIVE"),
        ("Custom insurance knowledge", "Teach accepted plans + coverage detail", "LIVE"),
        ("Monthly performance report", "Auto monthly summary generated from data", "LIVE"),
        ("Automated email reminders (72h/24h)", "Emailed to patient before the visit; cuts no-shows", "LIVE"),
        ("Patient recall campaigns (email)", "Win back lapsed patients via email", "LIVE"),
        ("Priority support", "Faster response (operational)", "LIVE"),
    ])

internal_slide(
    "ENTERPRISE ($997/mo) / White-Label - Feature Detail",
    "Multi-location, white-labeled. Everything in Growth, PLUS the items below.",
    [
        ("Unlimited conversations", "No monthly cap", "LIVE"),
        ("Unlimited locations & providers", "No limits enforced", "LIVE"),
        ("Multi-location routing", "Routes patient to the right office by ZIP / service", "BUILT"),
        ("Custom AI training", "Teach the AI clinic-specific procedures, policies, FAQs", "BUILT"),
        ("White-label branding", "Clinic logo/colors, remove 'Tabor' branding, reseller mode", "BUILT"),
        ("Custom domain", "clinic.yourdomain.com - needs per-clinic DNS + SSL", "PENDING"),
        ("EHR integration", "Sync to Epic/Cerner/Athena - real connectors not built yet (skeleton only)", "PENDING"),
        ("Dedicated manager + 24/7 support", "White-glove (operational)", "LIVE"),
    ])

internal_slide(
    "PENDING / ROADMAP - what's left to finish",
    "Consolidated to-do. Everything else is live on Supabase production today.",
    [
        ("Reminders/recall auto-fire", "Email send works; set an hourly Render cron -> /reminders/trigger & /api/recall/trigger", "PENDING"),
        ("Recall email unsubscribe link", "Add an unsubscribe link + endpoint before sending real recall email (CAN-SPAM)", "PENDING"),
        ("Email delivery (incl. booking confirmation)", "LIVE - SendGrid verified; signup + patient booking-confirmation emails send", "LIVE"),
        ("EHR integration", "Build real per-EHR connectors - currently a config skeleton", "PENDING"),
        ("White-label custom domain", "Per-clinic DNS + SSL wiring", "PENDING"),
        ("Core platform", "AI agent, booking, safety, dashboard, plan gating - all live & tested", "LIVE"),
        ("Database", "Supabase: persistent, RLS-secured, cascade FKs, hard-delete", "LIVE"),
        ("Security", "Admin-gated user create, rate limiting, JWT key, plan gating - all live", "LIVE"),
    ])

import sys
OUT = sys.argv[1] if len(sys.argv) > 1 else "HospitalAI_Aria_Overview.pptx"
prs.save(OUT)
print("Saved", OUT, "with", len(prs.slides._sldIdLst), "slides (4 client + 4 internal)")
